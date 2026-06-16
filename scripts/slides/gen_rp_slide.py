"""
gen_rp_slide.py — Generate the "Commercial Update" RP Meeting pipeline slide.

Stamps active-pipeline deal data into the branded Whitewater template, in the
same flow as page 2 of the Monthly Pipeline Report:  M3 -> M2.5 -> M2 -> M1.5 -> M1
(descending stage order), one clean line per deal, two balanced columns.

Why a template-stamping approach: the deck branding (black master, green stage
headers, Whitewater logo) lives in the .pptx itself. We edit the real template
rather than rebuilding branding in code, so the output is always pixel-faithful.

Usage:
    python gen_rp_slide.py --data deals.json [--template T.pptx] [--out O.pptx]
                           [--slide-index 1]

Input JSON schema (see scripts/slides/data/ for a worked example):
    {
      "title": "Commercial Update",
      "subtitle": "Increased lead generation ... Progress to M2 for key leads.",
      "deals": [
        {"stage":"M3","deal":"HP Inc - Global AI Marketing Reinvention - P1",
         "company":"HP","lead":"Nicolette Grams","amount":194700},
        ...
      ]
    }

Only `stage`, `company`, `amount` are strictly required per deal; `deal`
(full name) drives phase disambiguation when a company appears more than once,
and `lead` is the grey IE-lead column.
"""

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Brand palette (matches the existing deck) ────────────────────────────────
GREEN = RGBColor(0x68, 0xFF, 0xAE)   # stage headers
WHITE = RGBColor(0xFF, 0xFF, 0xFF)   # company + amount
GREY  = RGBColor(0xB5, 0xBC, 0xC6)   # IE lead + subtotals

# Descending pipeline order. Stages absent from the data are simply skipped.
STAGE_ORDER = ["M5", "M4", "M3.5", "M3", "M2.5", "M2", "M1.5", "M1", "M0"]

STAGE_FULL = {
    "M5":   "M5 — CONTRACT SIGNED",
    "M4":   "M4 — CONTRACT DRAFTING",
    "M3.5": "M3.5 — SOW REVISION",
    "M3":   "M3 — SOW PRESENTATION MEETING",
    "M2.5": "M2.5 — INTERNAL SOW DEVELOPMENT",
    "M2":   "M2 — CLIENT CO-DESIGN MEETING",
    "M1.5": "M1.5 — INTERNAL OPTIONS DEVELOPMENT",
    "M1":   "M1 — EXPLORATORY MEETING",
    "M0":   "M0 — LEAD",
}

# ── Layout geometry (inches) ─────────────────────────────────────────────────
START_Y   = 1.55
BOTTOM_Y  = 6.85          # keep clear of the footer at ~6.95"
LEFT_X    = 0.55
RIGHT_X   = 6.95
COL_W     = 5.85          # full column span for headers / subtotals
COMP_W    = 3.05          # company box
IE_DX     = 3.10          # company-x -> IE-x offset
IE_W      = 1.90
AMT_DX    = 5.05          # company-x -> amount-x offset
AMT_W     = 0.78

HDR_H, SUB_H, GAP_H = 0.30, 0.20, 0.12
DEAL_H_DEFAULT = 0.255    # auto-shrinks if a column would overflow

# Font sizes (pt)
HDR_PT, SUB_PT, COMP_PT, IE_PT, AMT_PT = 12.5, 10.5, 12.0, 10.5, 12.0


def fK(amount):
    """Format a dollar amount as the deck's compact $###.#K / $#.#M."""
    a = abs(amount)
    if a >= 1_000_000:
        return f"${amount / 1_000_000:.1f}M"
    return f"${amount / 1000:.1f}K"


def fTotal(amount):
    """Comma-grouped total, e.g. 1805296 -> $1,805.3K."""
    return f"${amount / 1000:,.1f}K"


PHASE_RE = re.compile(r"\bP(\d+)\b")
PHASE_WORD_RE = re.compile(r"\bPhase\s+(\d+)\b", re.IGNORECASE)


def deal_label(deal, repeated):
    """Primary bold label for a deal row.

    Company name is the clean default. When the same company appears more than
    once in the dataset we disambiguate by the phase token pulled from the deal
    name (P1/P2 or "Phase 3"); failing that we fall back to the full deal name
    so rows stay unique.
    """
    company = (deal.get("company") or "").strip()
    name = (deal.get("deal") or "").strip()
    if company and repeated.get(company, 0) > 1:
        m = PHASE_RE.search(name)
        if m:
            return f"{company} — P{m.group(1)}"
        m = PHASE_WORD_RE.search(name)
        if m:
            return f"{company} — Phase {m.group(1)}"
        return name or company
    return company or name


def group_by_stage(deals):
    """Return [(stage, [deals...]), ...] in descending pipeline order."""
    buckets = {}
    for d in deals:
        buckets.setdefault(d["stage"], []).append(d)
    ordered = []
    for sk in STAGE_ORDER:
        if sk in buckets:
            ordered.append((sk, buckets[sk]))
    # any unrecognised stage codes get appended at the end, in input order
    for sk in buckets:
        if sk not in STAGE_ORDER:
            ordered.append((sk, buckets[sk]))
    return ordered


def column_units(stages):
    """Approximate vertical 'units' a set of stages consumes, for balancing."""
    h = 0.0
    for _sk, ds in stages:
        h += HDR_H + SUB_H + len(ds) * DEAL_H_DEFAULT
    h += GAP_H * max(0, len(stages) - 1)
    return h


def split_columns(grouped):
    """Pick the contiguous left/right split that minimises the taller column."""
    if len(grouped) <= 1:
        return grouped, []
    best = None
    for cut in range(1, len(grouped)):
        left, right = grouped[:cut], grouped[cut:]
        taller = max(column_units(left), column_units(right))
        if best is None or taller < best[0]:
            best = (taller, left, right)
    return best[1], best[2]


def render_column(slide, items, col_x, deal_h, add_box):
    """Render one column's stages; returns the y-cursor after the last row."""
    comp_x = col_x
    ie_x = col_x + IE_DX
    amt_x = col_x + AMT_DX
    y = START_Y
    for i, (sk, ds) in enumerate(items):
        if i > 0:
            y += GAP_H
        add_box(comp_x, y, COL_W, HDR_H, STAGE_FULL.get(sk, sk), HDR_PT, True, GREEN)
        y += HDR_H
        total = sum(d["amount"] for d in ds)
        sub = f"{len(ds)} deal{'s' if len(ds) != 1 else ''} · {fTotal(total)}"
        add_box(comp_x, y, COL_W, SUB_H, sub, SUB_PT, False, GREY)
        y += SUB_H
        # company counts span the WHOLE dataset, injected by the caller
        repeated = render_column.repeated
        for d in ds:
            add_box(comp_x, y, COMP_W, deal_h, deal_label(d, repeated), COMP_PT, True, WHITE)
            lead = (d.get("lead") or "").strip()
            if lead:
                add_box(ie_x, y + 0.01, IE_W, deal_h, lead, IE_PT, False, GREY)
            add_box(amt_x, y, AMT_W, deal_h, fK(d["amount"]), AMT_PT, True, WHITE, PP_ALIGN.RIGHT)
            y += deal_h
    return y


def build(data, template_path, out_path, slide_index):
    prs = Presentation(str(template_path))
    if slide_index >= len(prs.slides):
        sys.exit(f"Template has {len(prs.slides)} slide(s); --slide-index {slide_index} out of range.")
    slide = prs.slides[slide_index]

    # 1) Strip the old deal content. The reference template holds the deal rows
    #    in a group named 'Group 6' (a scaled coordinate space — the source of
    #    every overlap bug) plus top-level 'Text N' boxes. Remove both; keep the
    #    title, subtitle and footers.
    for shape in list(slide.shapes):
        nm = shape.name or ""
        if nm == "Group 6" or re.match(r"^Text \d+$", nm):
            shape._element.getparent().remove(shape._element)

    # Optionally refresh the title / subtitle if provided.
    if data.get("title") or data.get("subtitle"):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            if data.get("title") and txt == "Commercial Update":
                _set_first_run(shape, data["title"])
            elif data.get("subtitle") and txt.lower().startswith("increased lead"):
                _set_first_run(shape, data["subtitle"])

    # 2) Helper to drop a clean, self-positioned text box.
    def add_box(left, top, width, height, text, size, bold, color, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.auto_size = None
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = "Arial"
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
        return tb

    deals = data["deals"]
    grouped = group_by_stage(deals)
    left, right = split_columns(grouped)

    # Whole-dataset company counts drive phase disambiguation.
    repeated = {}
    for d in deals:
        c = (d.get("company") or "").strip()
        repeated[c] = repeated.get(c, 0) + 1
    render_column.repeated = repeated

    # 3) Auto-shrink the deal row pitch if the taller column would overflow.
    deal_h = DEAL_H_DEFAULT
    for col in (left, right):
        rows = sum(len(ds) for _sk, ds in col)
        headers = len(col)
        fixed = headers * (HDR_H + SUB_H) + GAP_H * max(0, headers - 1)
        avail = BOTTOM_Y - START_Y - fixed
        if rows and avail / rows < deal_h:
            deal_h = max(0.16, avail / rows)

    end_l = render_column(slide, left, LEFT_X, deal_h, add_box)
    end_r = render_column(slide, right, RIGHT_X, deal_h, add_box) if right else START_Y

    prs.save(str(out_path))
    return {
        "deals": len(deals),
        "left_stages": [sk for sk, _ in left],
        "right_stages": [sk for sk, _ in right],
        "deal_h": round(deal_h, 3),
        "left_bottom": round(end_l, 2),
        "right_bottom": round(end_r, 2),
        "out": str(out_path),
    }


def _set_first_run(shape, text):
    p = shape.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for r in p.runs[1:]:
            r.text = ""
    else:
        p.text = text


def main():
    ap = argparse.ArgumentParser(description="Generate the RP Meeting pipeline slide.")
    ap.add_argument("--data", required=True, help="Path to deals JSON.")
    here = Path(__file__).resolve().parent
    ap.add_argument("--template", default=str(here / "templates" / "rp-meeting-slide.pptx"))
    ap.add_argument("--out", default=None, help="Output .pptx (default: alongside data).")
    ap.add_argument("--slide-index", type=int, default=1, help="0-based slide to stamp (default 1).")
    args = ap.parse_args()

    data = json.loads(Path(args.data).read_text(encoding="utf-8"))
    out = args.out or str(Path(args.data).with_suffix(".pptx"))
    info = build(data, args.template, out, args.slide_index)

    print(f"Stamped {info['deals']} deals -> {info['out']}")
    print(f"  Left column : {' '.join(info['left_stages'])}  (ends {info['left_bottom']}\")")
    print(f"  Right column: {' '.join(info['right_stages'])}  (ends {info['right_bottom']}\")")
    print(f"  Row pitch   : {info['deal_h']}\"  (footer clearance at {BOTTOM_Y}\")")


if __name__ == "__main__":
    main()
